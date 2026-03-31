import os
import json
from google import genai
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
import io
import chromadb
from chromadb.config import Settings
from typing import List, Dict, Any, Optional
from datetime import datetime
import uuid
import pandas as pd
import fitz # PyMuPDF
from docx import Document
from pptx import Presentation
from dotenv import load_dotenv

load_dotenv()

# Initialize Gemini Client
client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))

# Initialize ChromaDB (Local persistence)
chroma_client = chromadb.PersistentClient(path="./db_storage")
collection = chroma_client.get_or_create_collection(
    name="project_documents",
    metadata={"hnsw:space": "cosine"}
)

class DocumentProcessor:
    @staticmethod
    def parse_pdf(file_path: str) -> str:
        text = ""
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text()
        return text

    @staticmethod
    def parse_docx(file_path: str) -> str:
        doc = Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])

    @staticmethod
    def parse_pptx(file_path: str) -> str:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    @staticmethod
    def parse_xlsx_csv(file_path: str) -> str:
        if file_path.endswith('.xlsx'):
            df = pd.read_excel(file_path)
        else:
            df = pd.read_csv(file_path)
        return df.to_markdown(index=False)

    def process_file(self, file_path: str) -> str:
        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.pdf':
            return self.parse_pdf(file_path)
        elif ext == '.docx':
            return self.parse_docx(file_path)
        elif ext == '.pptx':
            return self.parse_pptx(file_path)
        elif ext in ['.xlsx', '.csv']:
            return self.parse_xlsx_csv(file_path)
        elif ext in ['.md', '.txt']:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        else:
            raise ValueError(f"Unsupported file format: {ext}")

class DriveService:
    def __init__(self):
        self.credentials_file = "google_credentials.json"
        self.folder_id = os.getenv("GOOGLE_DRIVE_FOLDER_ID")
        self.scopes = ['https://www.googleapis.com/auth/drive.readonly']
        self.service = None
        if os.path.exists(self.credentials_file):
            creds = service_account.Credentials.from_service_account_file(
                self.credentials_file, scopes=self.scopes)
            self.service = build('drive', 'v3', credentials=creds)

    def list_files(self, parent_id: Optional[str] = None) -> List[Dict[str, Any]]:
        # This now only fetches 1 level deep
        if not self.service:
            print("DriveService: No service initialized (credentials missing?)")
            return []
        
        target_folder = parent_id.replace("drive:", "") if parent_id else self.folder_id
        if not target_folder:
            print("DriveService: No target folder configured")
            return []

        try:
            query = f"'{target_folder}' in parents and trashed = false"
            results = self.service.files().list(
                q=query, fields="files(id, name, size, mimeType)").execute()
            
            drive_items = []
            for f in results.get('files', []):
                item_type = "folder" if f["mimeType"] == "application/vnd.google-apps.folder" else "file"
                drive_items.append({
                    "id": f"drive:{f['id']}",
                    "name": f["name"],
                    "type": item_type,
                    "size": f"{int(f.get('size', 0)) // 1024}KB" if item_type == "file" else None,
                    "children": [] if item_type == "folder" else None
                })
            return drive_items
        except Exception as e:
            print(f"DriveService Error listing files: {e}")
            return [] # Return empty list if drive fails, so local files can still show up

    def download_file(self, file_id: str) -> str:
        real_id = file_id.replace("drive:", "")
        meta = self.service.files().get(fileId=real_id, fields="id, name, mimeType").execute()
        
        if meta['mimeType'] == "application/vnd.google-apps.folder":
            return None

        filename = meta['name']
        local_path = os.path.join("docs_to_index", filename)
        
        # Handle Google Workspace Types
        is_workspace = meta['mimeType'].startswith('application/vnd.google-apps.')
        
        if is_workspace:
            local_path += ".pdf"
            fh = io.FileIO(local_path, 'wb')
            request = self.service.files().export_media(
                fileId=real_id, mimeType='application/pdf')
        else:
            fh = io.FileIO(local_path, 'wb')
            request = self.service.files().get_media(fileId=real_id)
            
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while done is False:
            status, done = downloader.next_chunk()
        return local_path

class SourceManagerService:
    def __init__(self):
        self.processor = DocumentProcessor()
        self.drive = DriveService()

    def _list_local(self, path: str) -> List[Dict[str, Any]]:
        items = []
        if not os.path.exists(path):
            return []
        for f in os.listdir(path):
            if f.startswith('.'): continue
            full_path = os.path.join(path, f)
            is_dir = os.path.isdir(full_path)
            items.append({
                "id": full_path,
                "name": f,
                "type": "folder" if is_dir else "file",
                "size": f"{os.path.getsize(full_path) // 1024}KB" if not is_dir else None,
                "children": [] if is_dir else None
            })
        return items

    async def list_available_files(self, parent_id: Optional[str] = None) -> List[Dict[str, Any]]:
        if not parent_id:
            # Top level: Local root + Drive root
            local = self._list_local("./docs_to_index")
            drive = self.drive.list_files()
            return local + drive
        
        # Level N: determine if local or drive
        if parent_id.startswith("drive:"):
            return self.drive.list_files(parent_id)
        else:
            return self._list_local(parent_id)

    async def sync_files(self, project_id: str, file_ids: List[str]) -> Dict[str, Any]:
        processed_count = 0
        processor = DocumentProcessor()
        
        if not os.path.exists("docs_to_index"):
            os.makedirs("docs_to_index")

        # Recursive indexing logic
        queue = list(file_ids)
        while queue:
            fid = queue.pop(0)
            
            # Case Local Folder: fetch children
            if os.path.isdir(fid):
                for f in os.listdir(fid):
                    if not f.startswith('.'):
                        queue.append(os.path.join(fid, f))
                continue

            try:
                # Case Drive ID
                if fid.startswith("drive:"):
                    real_id = fid.replace("drive:", "")
                    meta = self.drive.service.files().get(fileId=real_id, fields="id, name, mimeType").execute()
                    if meta['mimeType'] == "application/vnd.google-apps.folder":
                        # Fetch sub-files from Drive
                        children = self.drive.list_files(parent_id=fid)
                        for c in children:
                            queue.append(c['id'])
                        continue
                    else:
                        path = self.drive.download_file(fid)
                else:
                    path = fid
                
                # Parse & Index
                text = processor.process_file(path)
                if not text: continue
                
                chunks = [text[i:i + 1000] for i in range(0, len(text), 1000)]
                for i, chunk in enumerate(chunks):
                    embed_res = client.models.embed_content(
                        model='gemini-embedding-001',
                        contents=chunk,
                        config={'task_type': 'retrieval_document'}
                    )
                    collection.add(
                        ids=[f"{project_id}_{os.path.basename(path)}_{i}"],
                        embeddings=[embed_res.embeddings[0].values],
                        documents=[chunk],
                        metadatas=[{"project_id": project_id, "name": os.path.basename(path)}]
                    )
                processed_count += 1
            except Exception as e:
                print(f"Error indexing {fid}: {e}")

        return {"status": "success", "indexed_files": processed_count}

    async def purge_file(self, file_id: str):
        collection.delete(where={"file_id": file_id})
        return {"status": "deleted", "file_id": file_id}

    async def get_metrics(self) -> Dict[str, Any]:
        """Calculates storage size of ChromaDB and counts total embeddings."""
        import os
        
        def get_dir_size(path):
            total_size = 0
            if not os.path.exists(path):
                return 0
            for dirpath, dirnames, filenames in os.walk(path):
                for f in filenames:
                    fp = os.path.join(dirpath, f)
                    total_size += os.path.getsize(fp)
            return total_size

        db_size_bytes = get_dir_size("./db_storage")
        doc_count = collection.count()
        
        # Convert to readable format
        if db_size_bytes < 1024:
            size_str = f"{db_size_bytes} B"
        elif db_size_bytes < 1024 * 1024:
            size_str = f"{db_size_bytes / 1024:.2f} KB"
        else:
            size_str = f"{db_size_bytes / (1024 * 1024):.2f} MB"

        return {
            "storage_size": size_str,
            "total_chunks": doc_count,
            "status": "healthy"
        }

class ChatEngineService:
    def __init__(self):
        self.sessions_path = "db_storage/sessions.json"
        if not os.path.exists("db_storage"):
            os.makedirs("db_storage")
        if not os.path.exists(self.sessions_path):
            with open(self.sessions_path, 'w') as f:
                json.dump({}, f)

    def _load_sessions(self) -> Dict[str, Any]:
        try:
            with open(self.sessions_path, 'r') as f:
                return json.load(f)
        except:
            return {}

    def _save_sessions(self, sessions: Dict[str, Any]):
        with open(self.sessions_path, 'w') as f:
            json.dump(sessions, f, indent=4)

    async def list_sessions(self, project_id: str) -> List[Dict[str, Any]]:
        sessions = self._load_sessions()
        project_sessions = []
        for sid, sdata in sessions.items():
            if sdata.get("project_id") == project_id:
                project_sessions.append({
                    "id": sid,
                    "title": sdata.get("title", "Sans titre"),
                    "timestamp": sdata.get("created_at")
                })
        return sorted(project_sessions, key=lambda x: x['timestamp'], reverse=True)

    async def get_history(self, session_id: str) -> List[Dict[str, Any]]:
        sessions = self._load_sessions()
        if session_id in sessions:
            return sessions[session_id].get("messages", [])
        return []

    async def generate_response(self, query: str, project_id: str, session_id: Optional[str] = None) -> Dict[str, Any]:
        # 0. Load or Create Session
        sessions = self._load_sessions()
        if not session_id or session_id not in sessions:
            session_id = str(uuid.uuid4())
            sessions[session_id] = {
                "project_id": project_id,
                "title": query[:40] + "..." if len(query) > 40 else query,
                "created_at": datetime.now().isoformat(),
                "messages": []
            }
        
        # Save User Message
        user_msg = {
            "id": str(uuid.uuid4()),
            "role": "user",
            "content": query,
            "timestamp": datetime.now().isoformat()
        }
        sessions[session_id]["messages"].append(user_msg)

        # 1. Embed user query
        try:
            embed_res = client.models.embed_content(
                model='gemini-embedding-001',
                contents=query,
                config={'task_type': 'retrieval_query'}
            )
            query_embedding = embed_res.embeddings[0].values

            # 2. Search in ChromaDB
            results = collection.query(
                query_embeddings=[query_embedding],
                n_results=5,
                where={"project_id": project_id}
            )

            context = "\n---\n".join([doc for doc in results['documents'][0]]) if results['documents'] else ""
            source_names = list(set([met['name'] for met in results['metadatas'][0]])) if results['metadatas'] else []
        except Exception as e:
            print(f"RAG Retrieval Error: {e}")
            context = ""
            source_names = []

        # 3. Prompt
        system_instruction = (
            "Tu es ZeliveryBot, un assistant virtuel premium et expert dans l'analyse de documents projets. "
            "Ta mission est de répondre aux questions de l'utilisateur de manière claire, concise et professionnelle. "
            "Utilise impérativement le formatage Markdown pour rendre tes réponses lisibles : "
            "- Utilise des listes à puces pour énumérer des points. "
            "- Utilise le texte en gras (**texte**) pour souligner les concepts clés. "
            "- Utilise des titres (`### Titre`) pour structurer tes réponses longues. "
            "- Si une information pertinente est présente dans le [CONTEXTE] ci-dessous, utilise-la en priorité. "
            "Si le contexte est vide ou insuffisant, réponds de manière conviviale."
        )

        # Include Previous Messages (last 5) for better context
        history_str = ""
        for m in sessions[session_id]["messages"][-6:-1]: # Exclude current query which is already appended
            history_str += f"{m['role'].upper()}: {m['content']}\n"

        full_prompt = (
            f"{system_instruction}\n\n"
            f"[HISTORIQUE]\n{history_str}\n\n"
            f"[CONTEXTE]\n{context if context else 'Aucune donnée contextuelle trouvée.'}\n\n"
            f"[QUESTION]\n{query}"
        )

        # 4. Generate response
        response = client.models.generate_content(
            model='gemini-3-flash-preview',
            contents=full_prompt
        )

        ai_msg = {
            "id": str(uuid.uuid4()),
            "role": "assistant",
            "content": response.text,
            "sources": source_names,
            "timestamp": datetime.now().isoformat(),
            "session_id": session_id # Pass back the session ID
        }
        
        sessions[session_id]["messages"].append(ai_msg)
        self._save_sessions(sessions)

        return ai_msg

# Expose singleton services
source_manager = SourceManagerService()
chat_engine = ChatEngineService()
